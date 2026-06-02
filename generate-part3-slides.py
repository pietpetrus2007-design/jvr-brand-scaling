#!/usr/bin/env python3
"""Generate Part 3 PPTX slides for JvR Brand Scaling course."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

OUTPUT_DIR = "/tmp/part3-pptx"
os.makedirs(OUTPUT_DIR, exist_ok=True)

# Design constants
BLACK  = RGBColor(0x00, 0x00, 0x00)
ORANGE = RGBColor(0xFF, 0x6B, 0x00)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GRAY   = RGBColor(0xAA, 0xAA, 0xAA)

# 13.333" × 7.5" = standard 16:9 widescreen
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
FONT    = "Calibri"
LABEL   = "JvR Brand Scaling — Part 3"


# ── Helpers ──────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BLACK
    return slide


def rect(slide, x, y, w, h, color):
    s = slide.shapes.add_shape(1, x, y, w, h)  # 1 = RECTANGLE
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def txt(slide, text, x, y, w, h,
        size=28, bold=False, italic=False, color=WHITE, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size   = Pt(size)
    r.font.bold   = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name   = FONT
    return tb


def bullets(slide, items, x, y, w, h, size=26):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        if i > 0:
            p.space_before = Pt(10)
        # orange bullet
        r1 = p.add_run()
        r1.text = "●  "
        r1.font.size = Pt(size - 2)
        r1.font.color.rgb = ORANGE
        r1.font.name = FONT
        # white text
        r2 = p.add_run()
        r2.text = item
        r2.font.size = Pt(size)
        r2.font.color.rgb = WHITE
        r2.font.name = FONT
    return tb


# ── Slide types ───────────────────────────────────────────────────────────────

def title_slide(prs, title, label=LABEL):
    """Full-bleed title/cover slide."""
    slide = blank(prs)
    # Top + bottom bars
    rect(slide, 0, 0, SLIDE_W, Inches(0.10), ORANGE)
    rect(slide, 0, SLIDE_H - Inches(0.10), SLIDE_W, Inches(0.10), ORANGE)
    # Left vertical accent
    rect(slide, 0, 0, Inches(0.14), SLIDE_H, ORANGE)
    # Small decorative line above label
    rect(slide, Inches(0.55), Inches(0.95), Inches(0.70), Inches(0.06), ORANGE)
    # Course label
    txt(slide, label,
        Inches(0.55), Inches(1.15), Inches(12.4), Inches(0.55),
        size=20, color=GRAY)
    # Main title
    txt(slide, title,
        Inches(0.55), Inches(2.0), Inches(12.5), Inches(4.2),
        size=64, bold=True, color=ORANGE)
    return slide


def content_slide(prs, header, items):
    """Header + bullet-list slide."""
    slide = blank(prs)
    rect(slide, 0, 0, SLIDE_W, Inches(0.08), ORANGE)
    rect(slide, 0, 0, Inches(0.08), SLIDE_H, ORANGE)
    txt(slide, header,
        Inches(0.35), Inches(0.12), Inches(12.8), Inches(1.05),
        size=44, bold=True, color=ORANGE)
    rect(slide, Inches(0.35), Inches(1.22), Inches(12.8), Inches(0.05), ORANGE)
    bullets(slide, items,
            Inches(0.45), Inches(1.40), Inches(12.6), Inches(5.85))
    return slide


def rule_slide(prs, body):
    """'THE RULE' slide with large white text."""
    slide = blank(prs)
    rect(slide, 0, 0, SLIDE_W, Inches(0.08), ORANGE)
    rect(slide, 0, 0, Inches(0.08), SLIDE_H, ORANGE)
    txt(slide, "THE RULE",
        Inches(0.35), Inches(0.12), Inches(12.8), Inches(0.95),
        size=34, bold=True, color=ORANGE)
    rect(slide, Inches(0.35), Inches(1.1), Inches(12.8), Inches(0.05), ORANGE)
    txt(slide, body,
        Inches(0.45), Inches(1.28), Inches(12.5), Inches(5.9),
        size=34, color=WHITE)
    return slide


# ── Lesson content ────────────────────────────────────────────────────────────

LESSONS = {
    "p3m1l1": [
        ("title",   "Retainer vs Performance-Based"),
        ("content", "Two Ways to Get Paid", [
            "Retainer = flat monthly fee, same every month regardless of results",
            "Most common, most stable",
            "Start here",
        ]),
        ("content", "The Retainer", [
            "Flat monthly fee — you know exactly what is coming in",
            "Client knows exactly what they pay",
            "Try to keep above R2,000/month",
            "In the beginning — go lower if needed to get your first client",
            "Raise it as results come in",
        ]),
        ("content", "Performance-Based", [
            "% of revenue generated (5-15%)",
            "% of ad spend",
            "Per lead: R50-R200+ per lead, higher for high-ticket",
        ]),
        ("content", "Which One to Use?", [
            "Read the business",
            "Can afford you = retainer",
            "Hesitant to pay upfront but has potential = pilot (next lesson)",
            "Performance-based needs trust and clear tracking upfront",
        ]),
        ("rule",    "Retainer = stability.  Performance = upside.\n\n"
                    "Most healthy relationships start as retainer and evolve."),
    ],

    "p3m1l2": [
        ("title",   "How to Price Yourself"),
        ("content", "Start With Proof in Mind", [
            "Zero proof = R1,500-R2,000/month is fine",
            "First client is not about money — it is about your first result",
            "That result gets you client 2 at a higher rate",
        ]),
        ("content", "The R2,000 Floor", [
            "Once you have one result to show, never go below R2,000/month",
            "Aim for R2,000-R5,000 as your standard rate",
            "3 clients × R2,000 = R6,000/month",
        ]),
        ("content", "Moving Up", [
            "With multiple results, target R5,000-R15,000+ per client",
            "Bigger brands pay bigger fees",
            "You need proof to unlock this level",
            "Every client you land, your price goes up",
        ]),
        ("content", "Never Compete on Price", [
            "If a client says you are too expensive — do not drop your price",
            "Show more value or walk away",
            "Clients who negotiate your price from the start will be your worst clients",
        ]),
        ("rule",    "Charge what you would be embarrassed to ask for.\n\n"
                    "Then add 20%.\n\nThat is where you start."),
    ],

    "p3m1l3": [
        ("title",   "The Pilot Phase"),
        ("content", "What Is a Pilot?", [
            "A 2-week free trial offered to a specific type of client",
            "Not for everyone",
            "Only when a business has real potential and can afford you, "
            "but is too hesitant to commit upfront",
        ]),
        ("content", "When to Offer It", [
            "Business has budget and potential",
            "Owner is genuinely interested but nervous",
            "You believe you can get results quickly",
            "Do not offer a pilot to every client — it lowers your perceived value",
        ]),
        ("content", "How the Pilot Works", [
            "Run ads for 2 weeks at no charge",
            "Use their ad budget — you provide the service free",
            "At the end of 2 weeks, show results and present your retainer offer",
        ]),
        ("content", "Converting the Pilot", [
            "Show the results clearly after 2 weeks",
            "Present your monthly retainer as the next step",
            '"This is what we achieved in 2 weeks. Imagine what we do in 3 months."',
            "Do not ask — present",
        ]),
        ("rule",    "A pilot is not desperation.\n\n"
                    "It is a strategic move for the right client.\n\n"
                    "Used correctly it turns a hesitant business "
                    "into a long-term paying client."),
    ],

    "p3m2l1": [
        ("title",   "How to Invoice and Collect Payment"),
        ("content", "Two Ways to Collect", [
            "Invoice — use ChatGPT to generate a professional invoice, send as PDF",
            "Payment Link — Yoco, Stitch, or similar — client pays by card or EFT in 2 clicks",
        ]),
        ("content", "The Invoice", [
            'ChatGPT prompt: "Create a professional invoice for [your name/business], '
            'billing [client name] R[amount] for brand scaling services for [month]."',
            "Include invoice number, date, and payment due date of 7 days",
            "Export and send as PDF",
        ]),
        ("content", "The Payment Link", [
            "Create a payment link for the exact amount",
            "Send to client on WhatsApp or email",
            "Client pays in 2 clicks",
            "Money arrives in your account same day or next day",
        ]),
        ("content", "When to Send", [
            "Send invoice or payment link on the same day each month",
            "Make it consistent so clients expect it",
            "First of the month or last day of the month works well",
        ]),
        ("rule",    "Do not chase payment.\n\n"
                    "Send it, follow up once after 3 days if unpaid, then again at 7 days.\n\n"
                    "After that — address it directly."),
    ],

    "p3m2l2": [
        ("title",   "What to Do When a Client Does Not Pay"),
        ("content", "It Will Happen", [
            "At some point a client will delay or ignore payment",
            "Do not panic. Do not get emotional",
            "Handle it professionally and directly",
        ]),
        ("content", "Day 1-3", [
            "Send the invoice or payment link",
            "No follow-up yet — give them time",
        ]),
        ("content", "Day 4-7", [
            "One follow-up message:",
            '"Hey [name], just checking in on the invoice I sent. '
            'Let me know if you need anything from my side."',
        ]),
        ("content", "Day 7+", [
            "Call them directly or message clearly:",
            '"Hi [name], invoice is now 7 days overdue. I need this sorted before we '
            'continue with the campaigns. Let me know when payment goes through."',
        ]),
        ("content", "If They Still Do Not Pay", [
            "Pause the campaigns",
            "Do not threaten, do not beg",
            "Simply pause until payment is confirmed",
            "This usually solves it immediately",
        ]),
        ("rule",    "You are running a business, not a charity.\n\n"
                    "Be firm, be professional, but do not continue working for free."),
    ],

    "p3m2l3": [
        ("title",   "Asking for a Raise and Dropping Bad Clients"),
        ("content", "When to Ask for a Raise", [
            "You have been delivering results consistently",
            "It has been 2-3 months minimum",
            "You have proof to show",
            "The client is happy",
        ]),
        ("content", "How to Ask", [
            "Do not send a message saying 'I want to charge more'",
            '"I have been reviewing our results over the last few months. '
            'Based on what we have achieved I would like to adjust my rate '
            'to R[amount] from next month."',
            "Simple. Confident. No apology.",
        ]),
        ("content", "Bad Clients — Early Stage", [
            "In the beginning, every client teaches you something",
            "Take them. Use even small results as proof",
            "A bad client today is a stepping stone to a better one",
        ]),
        ("content", "When to Drop a Client", [
            "Once you have enough clients and proof, you can be selective",
            "Drop clients who do not pay on time consistently",
            "Drop clients who do not cooperate on content or approvals",
            "Drop clients who drain your time for little return",
        ]),
        ("rule",    "Every client you drop should be replaced with a better one.\n\n"
                    "Never drop without a plan.\n\n"
                    "Use bad clients to build proof. Then upgrade."),
    ],

    "p3m2l4": [
        ("title",   "What to Say When They Ask for Guarantees"),
        ("content", "They Will Ask", [
            '"Can you guarantee results?"',
            '"What if it does not work?"',
            "Know your answer before they ask it",
        ]),
        ("content", "The Honest Answer", [
            '"I cannot guarantee specific results because ads depend on your product, '
            'your market, and your budget."',
            '"What I can guarantee is that I will do everything in my power to get you '
            'the best possible return on your ad spend."',
        ]),
        ("content", "Back It Up", [
            "After that answer, show your process",
            "Explain how you test creatives, optimise campaigns, read data, and adjust",
            "Show them you know what you are doing",
            "Confidence is the real guarantee",
        ]),
        ("content", "Why Honesty Wins", [
            "Clients burned by fake promises will trust you more for being honest",
            "Those who still demand a guarantee after that are not the right clients",
        ]),
        ("rule",    "Never promise a number.\n\n"
                    "Promise your effort, your process, and your commitment.\n\n"
                    "Then deliver on it."),
    ],

    "p3m2l5": [
        ("title",   "The Mindset That Keeps You Going"),
        ("content", "This Is Not Easy", [
            "Brand scaling is one of the best business models available",
            "But it is not passive, it is not instant, and it will test your patience",
            "Know that going in",
        ]),
        ("content", "The First 90 Days Are the Hardest", [
            "Getting that first client takes time",
            "Getting results takes time. Building proof takes time",
            "Most people quit in this window",
            "The ones who push through are the ones who build something real",
        ]),
        ("content", "Every No Gets You Closer", [
            "Every outreach that does not reply is data",
            "Every call that does not close is data",
            "Every campaign that underperforms is data",
            "Do not treat it as failure — treat it as feedback",
        ]),
        ("content", "You Are Building a Skill", [
            "The difference between success and failure is not talent — it is repetition",
            "Every day you work on this you are getting better",
            "That compounds",
        ]),
        ("content", "The Standard", [
            "Show up every day",
            "Do the outreach. Run the ads. Learn from the results. Improve",
            "That is the entire game",
            "There is no secret beyond consistent execution",
        ]),
        ("rule",    "You now have everything you need.\n\n"
                    "The knowledge, the system, the tools.\n\n"
                    "The only thing left is you.\n\nGo build something real."),
    ],
}


# ── Build ─────────────────────────────────────────────────────────────────────

def build(lesson_id, slides):
    prs = new_prs()
    for s in slides:
        if s[0] == "title":
            title_slide(prs, s[1])
        elif s[0] == "content":
            content_slide(prs, s[1], s[2])
        elif s[0] == "rule":
            rule_slide(prs, s[1])
    path = os.path.join(OUTPUT_DIR, f"{lesson_id}.pptx")
    prs.save(path)
    print(f"  ✅  {lesson_id}  ({len(slides)} slides)  →  {path}")
    return len(slides)


if __name__ == "__main__":
    print(f"Generating Part 3 PPTX files → {OUTPUT_DIR}\n")
    counts = {}
    for lid, slides in LESSONS.items():
        counts[lid] = build(lid, slides)
    print("\nSlide counts:")
    for lid, n in counts.items():
        print(f"  {lid}: {n}")
    total = sum(counts.values())
    print(f"\n✅  Done — {len(counts)} lessons, {total} slides total")
